import os
import io
import re
import time
import base64
import json
import traceback
import uuid
import hashlib
from typing import Optional, Tuple, List, Dict, Any

from fastapi import FastAPI, UploadFile, File, Form, HTTPException
from fastapi.responses import JSONResponse, HTMLResponse, FileResponse
from fastapi.middleware.cors import CORSMiddleware

# Librairie Microsoft MarkItDown
from markitdown import MarkItDown

# Traitement d'images et PDF
from PIL import Image
import fitz  # PyMuPDF

# Traitement DOCX
import mammoth

# Traitement PPTX
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

# Traitement HTML
from html.parser import HTMLParser
import html

# ---------------------------
# Config via variables d'env
# ---------------------------
SAVE_UPLOADS  = os.getenv("SAVE_UPLOADS", "false").lower() == "true"
SAVE_OUTPUTS  = os.getenv("SAVE_OUTPUTS", "false").lower() == "true"
UPLOAD_DIR    = os.getenv("UPLOAD_DIR", "/data/uploads")
OUTPUT_DIR    = os.getenv("OUTPUT_DIR", "/data/outputs")
IMAGE_OUTPUT_DIR = os.getenv("IMAGE_OUTPUT_DIR", os.path.join(OUTPUT_DIR, "images"))
IMAGE_PUBLIC_BASE_URL = os.getenv("IMAGE_PUBLIC_BASE_URL", "").strip().rstrip("/")

# Images
IMG_FORMAT         = os.getenv("IMG_FORMAT", "png").strip().lower()
IMG_JPEG_QUALITY   = int(os.getenv("IMG_JPEG_QUALITY", "85"))
IMG_MAX_WIDTH      = int(os.getenv("IMG_MAX_WIDTH", "1400"))
IMG_ALT_PREFIX     = os.getenv("IMG_ALT_PREFIX", "Capture").strip()
MARKDOWN_INLINE_IMAGES = os.getenv("MARKDOWN_INLINE_IMAGES", "true").lower() in ("1", "true", "yes", "on")

# Dossiers persistents
os.makedirs(UPLOAD_DIR, exist_ok=True)
os.makedirs(OUTPUT_DIR, exist_ok=True)
os.makedirs(IMAGE_OUTPUT_DIR, exist_ok=True)

# ---------------------------
# App FastAPI
# ---------------------------
app = FastAPI(title="MarkItDown API", version="5.4-ui-ux")

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["POST", "GET"],
    allow_headers=["*"],
)

# ---------------------------
# Helpers génériques
# ---------------------------
def guess_file_type(filename: str, content_type: Optional[str]) -> str:
    fname = filename.lower()
    ctype = (content_type or "").lower()
    
    if ctype == "application/pdf" or fname.endswith(".pdf"):
        return "pdf"
    if "html" in ctype or fname.endswith((".html", ".htm")):
        return "html"
    if fname.endswith(".docx"):
        return "docx"
    if fname.endswith((".pptx", ".ppt")):
        return "pptx"
    if "json" in ctype or fname.endswith(".json"):
        return "json"
    if fname.endswith((".png", ".jpg", ".jpeg", ".bmp", ".tiff", ".webp")):
        return "image"
    return "unknown"

def _md_cleanup(md: str) -> str:
    """Nettoyage cosmétique du Markdown final."""
    if not md: return md
    lines = []
    for L in md.replace("\r", "").split("\n"):
        l = re.sub(r"[ \t]+$", "", L)
        l = re.sub(r"^\s*[•·●◦▪]\s+", "- ", l) # Normalisation puces
        l = re.sub(r"^\s*(\d+)[\)\]]\s+", r"\1. ", l) # Normalisation listes numérotées
        lines.append(l)
    return "\n".join(lines).strip()

def _safe_slug(value: str, fallback: str = "item") -> str:
    s = re.sub(r"[^a-zA-Z0-9._-]+", "-", value or "").strip("-._")
    return s[:120] or fallback

def _parse_bool(value: Optional[str], default: bool = False) -> bool:
    if value is None:
        return default
    return str(value).strip().lower() in ("1", "true", "yes", "on")

def _mime_ext(mime: str, fallback_fmt: str = "png") -> str:
    mime = (mime or "").lower()
    if "jpeg" in mime or "jpg" in mime:
        return "jpg"
    if "webp" in mime:
        return "webp"
    if "gif" in mime:
        return "gif"
    if "bmp" in mime:
        return "bmp"
    if "tiff" in mime:
        return "tiff"
    return "jpg" if fallback_fmt in ("jpg", "jpeg") else "png"

def _image_public_url(conversion_id: str, filename: str) -> str:
    rel = f"/images/{conversion_id}/{filename}"
    return f"{IMAGE_PUBLIC_BASE_URL}{rel}" if IMAGE_PUBLIC_BASE_URL else rel

def _encode_image_bytes(im: Image.Image, fmt: str = "png", quality: int = 85) -> Tuple[bytes, str, str]:
    buf = io.BytesIO()
    fmt = (fmt or "png").lower()
    if fmt in ("jpg", "jpeg"):
        im = im.convert("RGB")
        im.save(buf, format="JPEG", quality=quality, optimize=True)
        return buf.getvalue(), "image/jpeg", "jpg"
    im.save(buf, format="PNG", optimize=True)
    return buf.getvalue(), "image/png", "png"

def _save_image_asset(
    im: Image.Image,
    *,
    conversion_id: str,
    source_filename: str,
    image_index: int,
    page: Optional[int] = None,
    section_title: Optional[str] = None,
    caption: Optional[str] = None,
    context_before: str = "",
    context_after: str = "",
    bbox: Optional[Tuple[float, float, float, float]] = None,
) -> Dict[str, Any]:
    im = _pil_resize_max(im, IMG_MAX_WIDTH)
    data, mime, ext = _encode_image_bytes(im, IMG_FORMAT, IMG_JPEG_QUALITY)
    digest = hashlib.sha256(data).hexdigest()[:16]
    source_slug = _safe_slug(os.path.splitext(os.path.basename(source_filename or "document"))[0], "document")
    page_part = f"p{int(page):03d}" if page else "p000"
    image_id = f"{source_slug}_{page_part}_img{image_index:03d}_{digest}"
    filename = f"{image_id}.{ext}"
    out_dir = os.path.join(IMAGE_OUTPUT_DIR, _safe_slug(conversion_id, "conversion"))
    os.makedirs(out_dir, exist_ok=True)
    storage_path = os.path.join(out_dir, filename)
    with open(storage_path, "wb") as f:
        f.write(data)

    nearby = _normalize_text_spacing(" ".join([context_before or "", caption or "", context_after or ""]))
    return {
        "image_id": image_id,
        "image_index": image_index,
        "page": page,
        "section_title": section_title,
        "caption": caption or "",
        "context_before": context_before or "",
        "context_after": context_after or "",
        "context_nearby": nearby,
        "mime_type": mime,
        "width": im.width,
        "height": im.height,
        "storage_path": storage_path,
        "relative_url": f"/images/{conversion_id}/{filename}",
        "public_url": _image_public_url(conversion_id, filename),
        "bbox": list(bbox) if bbox else None,
        "content_sha256": digest,
    }

def _data_uri_to_image(data_uri: str) -> Optional[Tuple[Image.Image, str]]:
    if not data_uri:
        return None
    m = re.match(r"^data:(image/[^;]+);base64,(.+)$", data_uri, flags=re.IGNORECASE | re.DOTALL)
    if not m:
        return None
    try:
        raw = base64.b64decode(m.group(2), validate=False)
        return Image.open(io.BytesIO(raw)).copy(), m.group(1).lower()
    except Exception:
        return None

def _remove_headers_footers(md_lines: List[str]) -> List[str]:
    if len(md_lines) < 2: return md_lines
    first_lines, last_lines = [], []
    for c in md_lines:
        lines = [l for l in c.splitlines() if l.strip()]
        if lines:
            first_lines.append(lines[0])
            last_lines.append(lines[-1])
    
    # Détection répétitions (>50%)
    h_counts = {l: first_lines.count(l) for l in first_lines}
    f_counts = {l: last_lines.count(l) for l in last_lines}
    h_rem = {l for l, c in h_counts.items() if c >= 2 and c >= 0.5 * len(md_lines)}
    f_rem = {l for l, c in f_counts.items() if c >= 2 and c >= 0.5 * len(md_lines)}

    cleaned = []
    for c in md_lines:
        lines = c.splitlines()
        if lines and lines[0].strip() in h_rem: lines = lines[1:]
        if lines and lines[-1].strip() in f_rem: lines = lines[:-1]
        cleaned.append("\n".join(lines).strip())
    return cleaned

# --------------------------------------------------------------------
# Gestion Images Base64
# --------------------------------------------------------------------
def _pil_resize_max(im: Image.Image, max_w: int) -> Image.Image:
    if max_w and im.width > max_w:
        ratio = max_w / im.width
        new_h = int(im.height * ratio)
        return im.resize((max_w, new_h), Image.LANCZOS)
    return im

def _pil_to_base64(im: Image.Image, fmt: str = "png", quality: int = 85) -> str:
    buf = io.BytesIO()
    fmt = fmt.lower()
    if fmt in ("jpg", "jpeg"):
        im = im.convert("RGB")
        im.save(buf, format="JPEG", quality=quality, optimize=True)
        mime = "image/jpeg"
    else:
        im.save(buf, format="PNG", optimize=True)
        mime = "image/png"
    b64 = base64.b64encode(buf.getvalue()).decode("ascii")
    return f"data:{mime};base64,{b64}"

def _image_markdown(asset: Dict[str, Any], inline_images: bool, data_uri: Optional[str] = None) -> str:
    alt = asset.get("caption") or f"{IMG_ALT_PREFIX} {asset.get('image_index', '')}".strip()
    src = data_uri if inline_images and data_uri else asset.get("public_url") or asset.get("relative_url") or ""
    return f"![{alt}]({src})" if src else ""

def _crop_bbox_image(page: fitz.Page, bbox: Tuple[float, float, float, float], dpi: int = 200) -> Optional[Image.Image]:
    try:
        rect = fitz.Rect(*bbox)
        pix = page.get_pixmap(matrix=fitz.Matrix(dpi/72, dpi/72), clip=rect, alpha=False)
        return Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
    except: return None

# --------------------------------------------------------------------
# HTML / DOCX Helpers
# --------------------------------------------------------------------
_HTML_DATA_URI_IN_IMG = re.compile(r'<img\b[^>]*\bsrc=["\'](data:image/[^"\']+)["\']', flags=re.IGNORECASE)
_MD_DATA_IMG = re.compile(r'!\[[^\]]*\]\(\s*(?P<src>data:image/[^)]+)\s*\)')
_HTML_IMG_DATA_TAG = re.compile(r'<img\b[^>]*\bsrc=["\'](?P<src>data:image/[^"\']+)["\'][^>]*>', flags=re.IGNORECASE)

def _extract_html_data_imgs(html_text: str) -> List[str]:
    return _HTML_DATA_URI_IN_IMG.findall(html_text) if html_text else []

def _html_img_datauri_to_markdown(md: str, alt: str = "Capture") -> str:
    if not md or "<img" not in md.lower(): return md
    return _HTML_IMG_DATA_TAG.sub(lambda m: f'![{alt}]({m.group("src")})', md)

def _inject_full_data_uris_into_markdown(md: str, uris: List[str], prefix: str = "Capture") -> str:
    if not md or not uris: return md
    idx = 0
    def repl(m): nonlocal idx; src = uris[idx] if idx < len(uris) else m.group("src"); idx+=1; return f'![{prefix} – {idx}]({src})'
    return _MD_DATA_IMG.sub(repl, md)

# ---------------------------
# Renderers Spécifiques
# ---------------------------

def render_json_markdown(content: bytes) -> str:
    try:
        data = json.loads(content)
        if isinstance(data, list) and data and isinstance(data[0], dict):
            is_flat = True
            headers = []
            for row in data:
                if not isinstance(row, dict): continue
                for k in row.keys():
                    if k not in headers: headers.append(k)
                for v in row.values():
                    if isinstance(v, (dict, list)):
                        is_flat = False
                        break
                if not is_flat: break
            
            if is_flat and headers:
                lines = []
                lines.append("| " + " | ".join(str(h) for h in headers) + " |")
                lines.append("| " + " | ".join(["---"] * len(headers)) + " |")
                for row in data:
                    if not isinstance(row, dict): continue
                    row_vals = []
                    for h in headers:
                        val = row.get(h, "")
                        s_val = str(val).replace("|", "&#124;").replace("\n", "<br>")
                        row_vals.append(s_val)
                    lines.append("| " + " | ".join(row_vals) + " |")
                return "\n".join(lines)

        elif isinstance(data, dict):
            is_flat = all(not isinstance(v, (dict, list)) for v in data.values())
            if is_flat and data:
                lines = ["| Clé | Valeur |", "| --- | --- |"]
                for k, v in data.items():
                    s_val = str(v).replace("|", "&#124;").replace("\n", "<br>")
                    lines.append(f"| {k} | {s_val} |")
                return "\n".join(lines)

        formatted = json.dumps(data, indent=2, ensure_ascii=False)
        return f"```json\n{formatted}\n```"
    except json.JSONDecodeError:
        return "```text\n(Fichier JSON invalide)\n```"

def render_pptx_markdown(content: bytes, source_filename: str, conversion_id: str, inline_images: bool = False) -> Tuple[str, List[Dict[str, Any]]]:
    try:
        prs = Presentation(io.BytesIO(content))
        md_slides = []
        images: List[Dict[str, Any]] = []
        image_index = 0
        for i, slide in enumerate(prs.slides):
            slide_content = []
            if slide.shapes.title:
                slide_content.append(f"## {slide.shapes.title.text}")
            else:
                slide_content.append(f"## Slide {i+1}")
            shapes = list(slide.shapes)
            shapes.sort(key=lambda s: (s.top, s.left))
            for shape in shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        txt = paragraph.text.strip()
                        if txt:
                            prefix = "- " if paragraph.level > 0 else ""
                            slide_content.append(f"{prefix}{txt}")
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    try:
                        image_blob = shape.image.blob
                        with Image.open(io.BytesIO(image_blob)) as im:
                            image_index += 1
                            caption = f"{IMG_ALT_PREFIX} - slide {i+1}"
                            before = "\n".join(slide_content[-3:])
                            asset = _save_image_asset(
                                im,
                                conversion_id=conversion_id,
                                source_filename=source_filename,
                                image_index=image_index,
                                page=i + 1,
                                section_title=slide.shapes.title.text.strip() if slide.shapes.title else None,
                                caption=caption,
                                context_before=before,
                            )
                            images.append(asset)
                            if inline_images:
                                data_uri = _pil_to_base64(im, IMG_FORMAT, IMG_JPEG_QUALITY)
                                slide_content.append(_image_markdown(asset, True, data_uri))
                    except Exception: pass
            md_slides.append("\n\n".join(slide_content))
        return "\n\n---\n\n".join(md_slides), images
    except Exception as e:
        return f"Erreur PPTX: {str(e)}", []

def render_docx_markdown(content: bytes, source_filename: str, conversion_id: str, inline_images: bool = False) -> Tuple[str, Dict, List[Dict[str, Any]]]:
    try:
        res = mammoth.convert_to_html(io.BytesIO(content))
        uris = _extract_html_data_imgs(res.value)
        images: List[Dict[str, Any]] = []
        replacement_uris: List[str] = []
        for idx, uri in enumerate(uris, start=1):
            decoded = _data_uri_to_image(uri)
            if not decoded:
                continue
            im, _mime = decoded
            asset = _save_image_asset(
                im,
                conversion_id=conversion_id,
                source_filename=source_filename,
                image_index=idx,
                caption=f"{IMG_ALT_PREFIX} - image {idx}",
            )
            images.append(asset)
            replacement_uris.append(uri if inline_images else asset.get("public_url") or asset.get("relative_url") or "")

        md_eng = MarkItDown()
        out = md_eng.convert_stream(io.BytesIO(res.value.encode('utf-8')), file_extension=".html")
        md = getattr(out, "text_content", "") or ""
        md = _html_img_datauri_to_markdown(md)
        md = _inject_full_data_uris_into_markdown(md, replacement_uris, IMG_ALT_PREFIX)
        if not inline_images:
            md = re.sub(r'!\[[^\]]*\]\([^)]+\)', "", md)
        return _md_cleanup(md), {"messages": [m.message for m in res.messages]}, images
    except Exception as e: return f"Erreur DOCX: {e}", {"error": str(e)}, []

def _normalize_text_spacing(text: str) -> str:
    if not text:
        return ""
    text = re.sub(r"[ 	]+", " ", text)
    text = re.sub(r"\s+([,.;:!?])", r"\1", text)
    text = re.sub(r"([\[(])\s+", r"\1", text)
    text = re.sub(r"\s+([\])])", r"\1", text)
    return text.strip()


def _extract_pdf_text_blocks(raw: Dict[str, Any], median: float) -> List[Dict[str, Any]]:
    blocks = []
    for b in raw.get("blocks", []):
        if b.get("type") != 0:
            continue
        bbox = b.get("bbox")
        if not bbox:
            continue

        block_lines = []
        max_s = 0.0
        is_bold = False
        lines = sorted(
            b.get("lines", []),
            key=lambda l: (
                round(((l.get("bbox") or [0, 0, 0, 0])[1]) / 2),
                round((l.get("bbox") or [0, 0, 0, 0])[0], 1),
            ),
        )
        for line in lines:
            spans = sorted(
                line.get("spans", []),
                key=lambda s: (
                    round((s.get("bbox") or [0, 0, 0, 0])[0], 1),
                    round((s.get("bbox") or [0, 0, 0, 0])[1], 1),
                ),
            )
            line_parts = []
            prev_x1 = None
            for s in spans:
                t = (s.get("text") or "")
                if not t.strip():
                    continue
                s_bbox = s.get("bbox") or [0, 0, 0, 0]
                x0, x1 = s_bbox[0], s_bbox[2]
                if prev_x1 is not None and x0 - prev_x1 > max(3.0, float(s.get("size", 0.0)) * 0.15):
                    line_parts.append(" ")
                max_s = max(max_s, float(s.get("size", 0.0)))
                flags = int(s.get("flags", 0))
                is_bold = is_bold or _is_bold(flags)
                line_parts.append(f"**{t}**" if _is_bold(flags) else t)
                prev_x1 = x1
            line_text = _normalize_text_spacing("".join(line_parts))
            if line_text:
                block_lines.append(line_text)

        full = _normalize_text_spacing(" ".join(block_lines))
        if not full:
            continue

        h = _classify_heading(max_s, median, is_bold)
        blocks.append({
            "type": "text",
            "x0": float(bbox[0]),
            "y0": float(bbox[1]),
            "x1": float(bbox[2]),
            "y1": float(bbox[3]),
            "text": f"{h} {full}" if h else full,
        })
    return blocks


def _horizontal_overlap_ratio(a: Dict[str, Any], b: Dict[str, Any]) -> float:
    overlap = max(0.0, min(a["x1"], b["x1"]) - max(a["x0"], b["x0"]))
    width = max(1.0, min(a["x1"] - a["x0"], b["x1"] - b["x0"]))
    return overlap / width


def _same_column(a: Dict[str, Any], b: Dict[str, Any], x_tol: float = 36.0, overlap_ratio: float = 0.35) -> bool:
    a_center = (a["x0"] + a["x1"]) / 2.0
    b_center = (b["x0"] + b["x1"]) / 2.0
    left_close = abs(a["x0"] - b["x0"]) <= x_tol
    center_close = abs(a_center - b_center) <= max(x_tol, min(a["x1"] - a["x0"], b["x1"] - b["x0"]) * 0.2)
    return left_close or center_close or _horizontal_overlap_ratio(a, b) >= overlap_ratio


def _vertical_gap(a: Dict[str, Any], b: Dict[str, Any]) -> float:
    return max(0.0, b["y0"] - a["y1"])


def _find_pdf_target_group(groups: List[Dict[str, Any]], block: Dict[str, Any]) -> Optional[int]:
    best_idx: Optional[int] = None
    best_score: Optional[Tuple[float, float, float]] = None

    for idx, group in enumerate(groups):
        if not _same_column(group, block):
            continue

        gap = block["y0"] - group["y1"]
        if gap < -2.0:
            continue

        max_gap = max(26.0, min(group["y1"] - group["y0"], block["y1"] - block["y0"]) * 1.2)
        if gap > max_gap:
            continue

        overlap = _horizontal_overlap_ratio(group, block)
        center_gap = abs(((group["x0"] + group["x1"]) / 2.0) - ((block["x0"] + block["x1"]) / 2.0))
        score = (gap, center_gap, -overlap)
        if best_score is None or score < best_score:
            best_score = score
            best_idx = idx

    return best_idx


def _merge_pdf_text_blocks(blocks: List[Dict[str, Any]]) -> List[Dict[str, Any]]:
    if not blocks:
        return []

    ordered = sorted(blocks, key=lambda b: (round(b["y0"] / 4), b["x0"], b["y1"], b["x1"]))
    groups: List[Dict[str, Any]] = []

    for block in ordered:
        idx = _find_pdf_target_group(groups, block)
        if idx is None:
            groups.append(dict(block))
            continue

        group = groups[idx]
        group["text"] = _normalize_text_spacing(f"{group['text']} {block['text']}")
        group["x0"] = min(group["x0"], block["x0"])
        group["y0"] = min(group["y0"], block["y0"])
        group["x1"] = max(group["x1"], block["x1"])
        group["y1"] = max(group["y1"], block["y1"])

    return sorted(groups, key=lambda g: (round(g["y0"] / 6), g["x0"], g["y1"], g["x1"]))


def _extract_pdf_image_blocks(
    page: fitz.Page,
    raw: Dict[str, Any],
    page_index: int,
    *,
    source_filename: str,
    conversion_id: str,
    start_index: int = 0,
) -> Tuple[List[Dict[str, Any]], List[Dict[str, Any]]]:
    items = []
    images: List[Dict[str, Any]] = []
    image_index = start_index
    for b in raw.get("blocks", []):
        if b.get("type") != 1:
            continue
        bbox = b.get("bbox")
        if not bbox:
            continue
        im = _crop_bbox_image(page, bbox, 300)
        if not im:
            continue
        image_index += 1
        asset = _save_image_asset(
            im,
            conversion_id=conversion_id,
            source_filename=source_filename,
            image_index=image_index,
            page=page_index,
            caption=f"{IMG_ALT_PREFIX} - page {page_index}",
            bbox=tuple(float(x) for x in bbox),
        )
        images.append(asset)
        items.append({
            "type": "image",
            "x0": float(bbox[0]),
            "y0": float(bbox[1]),
            "x1": float(bbox[2]),
            "y1": float(bbox[3]),
            "asset": asset,
            "text": "",
        })
    return items, images


def _current_section_from_text(text: str, current: Optional[str]) -> Optional[str]:
    for line in (text or "").splitlines():
        stripped = line.strip()
        if stripped.startswith("#"):
            return stripped.lstrip("#").strip() or current
    return current


def _annotate_pdf_image_context(items: List[Dict[str, Any]], inline_images: bool) -> None:
    current_section: Optional[str] = None
    for idx, item in enumerate(items):
        if item.get("type") == "text":
            current_section = _current_section_from_text(item.get("text") or "", current_section)
            continue
        if item.get("type") != "image" or not item.get("asset"):
            continue

        before_parts = []
        for prev in reversed(items[:idx]):
            if prev.get("type") == "text" and prev.get("text"):
                before_parts.append(prev["text"])
            if len(before_parts) >= 2:
                break
        before = _normalize_text_spacing(" ".join(reversed(before_parts)))[:1200]

        after_parts = []
        for nxt in items[idx + 1:]:
            if nxt.get("type") == "text" and nxt.get("text"):
                after_parts.append(nxt["text"])
            if len(after_parts) >= 2:
                break
        after = _normalize_text_spacing(" ".join(after_parts))[:1200]

        asset = item["asset"]
        asset["section_title"] = asset.get("section_title") or current_section
        asset["context_before"] = before
        asset["context_after"] = after
        asset["context_nearby"] = _normalize_text_spacing(" ".join([before, asset.get("caption") or "", after]))
        item["text"] = _image_markdown(asset, inline_images) if inline_images else ""


def render_pdf_markdown_inline(
    content: bytes,
    meta: Dict,
    source_filename: str,
    conversion_id: str,
    inline_images: bool = False,
) -> Tuple[str, List[Dict[str, Any]]]:
    doc = fitz.open(stream=content, filetype="pdf")
    meta["pages"] = doc.page_count
    md_pages = []
    images: List[Dict[str, Any]] = []
    try:
        for p in range(doc.page_count):
            page = doc.load_page(p)
            raw = page.get_text("dict", sort=False) or {}
            median = _median_font_size(raw)

            text_blocks = _merge_pdf_text_blocks(_extract_pdf_text_blocks(raw, median))
            image_blocks, page_images = _extract_pdf_image_blocks(
                page,
                raw,
                p + 1,
                source_filename=source_filename,
                conversion_id=conversion_id,
                start_index=len(images),
            )
            images.extend(page_images)
            items = sorted(
                text_blocks + image_blocks,
                key=lambda item: (round(item["y0"] / 6), item["x0"], item["y1"], item["x1"]),
            )
            _annotate_pdf_image_context(items, inline_images)
            md_pages.append("\n\n".join(item["text"] for item in items if item.get("text")))

        cleaned = _remove_headers_footers(md_pages)
        return _md_cleanup("\n\n---\n\n".join(cleaned)), images
    finally:
        doc.close()

# PDF Utils
def _is_bold(flags): return bool(flags & 1 or flags & 32)
def _median_font_size(raw):
    sizes = [s["size"] for b in raw.get("blocks",[]) if b["type"]==0 for l in b["lines"] for s in l["spans"] if s["text"].strip()]
    if not sizes: return 0.0
    sizes.sort(); mid = len(sizes)//2
    return sizes[mid] if len(sizes)%2 else (sizes[mid-1]+sizes[mid])/2.0
def _classify_heading(s, med, b):
    if med <= 0: return None
    if s >= med*1.8: return "#"
    if s >= med*1.5: return "##"
    if b and s >= med*1.1: return "###"
    return None

# ---------------------------
# UI
# ---------------------------
HTML_PAGE = r'''<!doctype html>
<html lang="fr">
<head>
  <meta charset="utf-8"/>
  <meta name="viewport" content="width=device-width,initial-scale=1"/>
  <title>MarkItDown Converter</title>
  <script src="https://cdn.jsdelivr.net/npm/marked/marked.min.js"></script>
  <link rel="stylesheet" href="https://cdnjs.cloudflare.com/ajax/libs/font-awesome/6.0.0/css/all.min.css">
  <style>
    :root{
      color-scheme: dark;
      --bg: #0f172a; --bg-card: #1e293b; --bg-input: #334155;
      --text: #f1f5f9; --text-muted: #94a3b8;
      --primary: #3b82f6; --primary-hover: #2563eb;
      --border: #334155;
      --success: #10b981;
    }
    *{box-sizing:border-box}
    body{ font-family: 'Inter', -apple-system, sans-serif; background: var(--bg); color: var(--text); margin: 0; display:flex; flex-direction:column; height:100vh; overflow:hidden; }
    
    /* Header */
    header { height: 60px; background: var(--bg-card); border-bottom: 1px solid var(--border); display: flex; align-items: center; padding: 0 1.5rem; justify-content: space-between; }
    .brand { display:flex; align-items:center; gap:0.75rem; font-weight:700; font-size:1.1rem; color: #fff; }
    .brand i { color: var(--primary); font-size:1.2rem; }
    .actions { display: flex; gap: 0.5rem; }

    /* Main Layout */
    main { flex: 1; padding: 1rem; display: flex; gap: 1rem; overflow: hidden; position: relative; }
    
    .sidebar { width: 320px; display: flex; flex-direction: column; gap: 1rem; flex-shrink: 0; }
    .content-area { flex: 1; display: flex; gap: 1rem; overflow: hidden; transition: all 0.3s ease; }
    
    /* Cards & Dropzone */
    .card { background: var(--bg-card); border: 1px solid var(--border); border-radius: 8px; padding: 1rem; display: flex; flex-direction: column; }
    
    .drop-zone {
        border: 2px dashed var(--border); border-radius: 8px; padding: 1.5rem 1rem;
        text-align: center; color: var(--text-muted); cursor: pointer; transition: 0.2s;
        display: flex; flex-direction: column; align-items: center; gap: 0.5rem;
    }
    .drop-zone:hover, .drop-zone.active { border-color: var(--primary); background: rgba(59, 130, 246, 0.05); }
    .drop-zone i { font-size: 1.5rem; color: var(--text-muted); margin-bottom: 0.25rem; }
    
    .file-info { margin-top: 1rem; background: var(--bg-input); padding: 0.5rem; border-radius: 6px; font-size: 0.85rem; display: flex; align-items: center; justify-content: space-between; }
    .file-info span { overflow: hidden; text-overflow: ellipsis; white-space: nowrap; max-width: 200px; }
    
    /* Stats */
    .stats-row { display: flex; gap: 0.5rem; margin-top: 0.5rem; flex-wrap: wrap; }
    .stat-badge { background: rgba(255,255,255,0.05); padding: 2px 8px; border-radius: 99px; font-size: 0.75rem; color: var(--text-muted); }

    /* Editors */
    .editor-container, .preview-container {
        flex: 1; display: flex; flex-direction: column; background: var(--bg-card);
        border: 1px solid var(--border); border-radius: 8px; overflow: hidden; position: relative;
    }
    
    .panel-header {
        padding: 0.5rem 1rem; border-bottom: 1px solid var(--border); background: rgba(0,0,0,0.1);
        display: flex; justify-content: space-between; align-items: center; font-size: 0.85rem; font-weight: 600; color: var(--text-muted);
    }
    
    .toolbar-btn {
        background: transparent; border: none; color: var(--text-muted); cursor: pointer;
        padding: 4px; border-radius: 4px; transition: 0.2s;
    }
    .toolbar-btn:hover { color: var(--text); background: rgba(255,255,255,0.1); }

    textarea, #preview-content {
        flex: 1; width: 100%; border: none; background: #0f1319; color: #e2e8f0;
        padding: 1rem; font-family: 'Menlo', monospace; font-size: 0.9rem; resize: none; outline: none;
        overflow-y: auto; line-height: 1.6;
    }
    
    /* Preview Styling */
    #preview-content { font-family: -apple-system, sans-serif; background: var(--bg-card); }
    #preview-content img { max-width: 100%; border-radius: 4px; border: 1px solid var(--border); margin: 0.5rem 0; }
    #preview-content h1, #preview-content h2 { border-bottom: 1px solid var(--border); padding-bottom: 0.3em; margin-top: 1.5em; color: #fff; }
    #preview-content pre { background: #0f1319; padding: 1rem; border-radius: 6px; overflow-x: auto; }
    #preview-content code { font-family: 'Menlo', monospace; font-size: 0.85em; background: rgba(255,255,255,0.05); padding: 0.2em 0.4em; border-radius: 3px; }
    #preview-content pre code { background: transparent; padding: 0; }
    #preview-content table { width: 100%; border-collapse: collapse; margin: 1em 0; }
    #preview-content th, #preview-content td { border: 1px solid var(--border); padding: 8px 12px; }
    #preview-content th { background: rgba(0,0,0,0.2); }
    #preview-content blockquote { border-left: 3px solid var(--primary); margin: 1em 0; padding-left: 1em; color: var(--text-muted); }

    /* Buttons */
    .btn {
        padding: 0.6rem 1rem; border-radius: 6px; border: none; font-weight: 600; cursor: pointer;
        display: inline-flex; align-items: center; gap: 0.5rem; transition: 0.2s; font-size: 0.9rem;
    }
    .btn-primary { background: var(--primary); color: white; width: 100%; justify-content: center; margin-top: 1rem; }
    .btn-primary:hover:not(:disabled) { background: var(--primary-hover); }
    .btn-primary:disabled { opacity: 0.6; cursor: wait; }
    
    .toggle-preview { display: flex; align-items: center; gap: 0.4rem; font-size: 0.85rem; color: var(--text-muted); cursor: pointer; user-select: none; }
    .toggle-preview:hover { color: var(--text); }
    
    /* Utilities */
    .hidden { display: none !important; }
    
    /* Toast Notification */
    .toast {
        position: fixed; bottom: 20px; right: 20px; background: var(--bg-card); border: 1px solid var(--border);
        padding: 12px 20px; border-radius: 8px; box-shadow: 0 4px 12px rgba(0,0,0,0.3);
        display: flex; align-items: center; gap: 10px; transform: translateY(100px); transition: transform 0.3s cubic-bezier(0.175, 0.885, 0.32, 1.275);
        z-index: 9999;
    }
    .toast.show { transform: translateY(0); }
    .toast i { color: var(--success); }

    /* Responsive */
    @media (max-width: 900px) {
        main { flex-direction: column; overflow-y: auto; }
        .sidebar { width: 100%; }
        .content-area { flex-direction: column; overflow: visible; height: auto; }
        .editor-container, .preview-container { height: 500px; flex: none; }
    }
  </style>
</head>
<body>

  <header>
    <div class="brand">
      <i class="fas fa-file-invoice"></i>
      <span>MarkItDown</span>
    </div>
    <div class="actions">
       <div class="toggle-preview" id="toggleViewBtn" onclick="togglePreview()">
         <i class="fas fa-columns"></i> <span>Aperçu</span>
       </div>
    </div>
  </header>

  <main>
    <aside class="sidebar">
      <div class="card">
        <div class="drop-zone" id="dropzone">
          <i class="fas fa-cloud-upload-alt"></i>
          <span>Glisser-déposer un fichier<br><small>(PDF, DOCX, PPTX, JSON, IMG)</small></span>
        </div>
        <input type="file" id="fileInput" hidden />
        
        <div id="fileInfo" class="file-info hidden">
          <span id="fileName">fichier.pdf</span>
          <i class="fas fa-times" style="cursor:pointer" onclick="resetFile()"></i>
        </div>

        <button id="convertBtn" class="btn btn-primary" disabled onclick="convertFile()">
          <i class="fas fa-bolt"></i> Convertir
        </button>

        <div class="stats-row">
           <span class="stat-badge hidden" id="statusBadge">Prêt</span>
           <span class="stat-badge hidden" id="timerBadge">0s</span>
        </div>
      </div>
      
      <div class="card" style="flex:1; justify-content:center; align-items:center; text-align:center; opacity:0.5">
        <small>Powered by Yoslo59<br>&copy; 2025-2026</small>
      </div>
    </aside>

    <div class="content-area" id="contentArea">
      <!-- Editor -->
      <div class="editor-container">
        <div class="panel-header">
          <span>MARKDOWN</span>
          <div style="display:flex; gap:0.5rem">
            <button class="toolbar-btn" onclick="copyToClipboard()" title="Copier"><i class="far fa-copy"></i></button>
            <button class="toolbar-btn" id="downloadBtn" onclick="downloadFile()" title="Télécharger" disabled><i class="fas fa-download"></i></button>
          </div>
        </div>
        <textarea id="markdownOutput" spellcheck="false" placeholder="Le résultat apparaîtra ici..."></textarea>
      </div>

      <!-- Preview -->
      <div class="preview-container" id="previewPanel">
        <div class="panel-header">
          <span>APERÇU</span>
        </div>
        <div id="preview-content">
            <div style="height:100%; display:flex; align-items:center; justify-content:center; color:var(--text-muted);">
                Aperçu du rendu
            </div>
        </div>
      </div>
    </div>
  </main>

  <!-- Toast -->
  <div class="toast" id="toast">
    <i class="fas fa-check-circle"></i>
    <span id="toastMsg">Action effectuée</span>
  </div>

<script>
const $ = id => document.getElementById(id);
let currentBlobUrl = null;
let isPreviewVisible = true;

// Init
marked.use({ breaks: true, gfm: true });
const mdOut = $("markdownOutput");
const prevContent = $("preview-content");

// --- UI Logic ---

function togglePreview() {
    const panel = $("previewPanel");
    const btn = $("toggleViewBtn");
    isPreviewVisible = !isPreviewVisible;
    
    if (isPreviewVisible) {
        panel.classList.remove("hidden");
        btn.style.opacity = "1";
        btn.querySelector("span").innerText = "Aperçu";
    } else {
        panel.classList.add("hidden");
        btn.style.opacity = "0.6";
        btn.querySelector("span").innerText = "Code seul";
    }
}

function updatePreview() {
    const text = mdOut.value;
    if (!text.trim()) {
        prevContent.innerHTML = '<div style="height:100%; display:flex; align-items:center; justify-content:center; color:var(--text-muted);">Aperçu du rendu</div>';
    } else {
        prevContent.innerHTML = marked.parse(text);
    }
}

function showToast(msg, type="success") {
    const t = $("toast");
    $("toastMsg").innerText = msg;
    t.querySelector("i").className = type === "error" ? "fas fa-exclamation-circle" : "fas fa-check-circle";
    t.querySelector("i").style.color = type === "error" ? "#ef4444" : "#10b981";
    t.classList.add("show");
    setTimeout(() => t.classList.remove("show"), 3000);
}

// --- File Logic ---

$("dropzone").onclick = () => $("fileInput").click();
$("dropzone").ondragover = e => { e.preventDefault(); $("dropzone").classList.add("active"); };
$("dropzone").ondragleave = () => $("dropzone").classList.remove("active");
$("dropzone").ondrop = e => { 
    e.preventDefault(); 
    $("dropzone").classList.remove("active"); 
    if(e.dataTransfer.files[0]) handleFile(e.dataTransfer.files[0]); 
};
$("fileInput").onchange = () => { if($("fileInput").files[0]) handleFile($("fileInput").files[0]); };

function handleFile(f) {
    // Hack pour réassigner le fichier au input si drop
    const dt = new DataTransfer(); dt.items.add(f); $("fileInput").files = dt.files;
    
    $("fileName").innerText = f.name;
    $("fileInfo").classList.remove("hidden");
    $("dropzone").classList.add("hidden");
    $("convertBtn").disabled = false;
    $("statusBadge").classList.add("hidden");
}

function resetFile() {
    $("fileInput").value = "";
    $("fileInfo").classList.add("hidden");
    $("dropzone").classList.remove("hidden");
    $("convertBtn").disabled = true;
    mdOut.value = "";
    updatePreview();
}

async function convertFile() {
    const f = $("fileInput").files[0];
    if(!f) return;
    
    const btn = $("convertBtn");
    btn.disabled = true;
    btn.innerHTML = '<i class="fas fa-circle-notch fa-spin"></i> Traitement...';
    $("statusBadge").classList.remove("hidden");
    $("statusBadge").innerText = "Conversion...";
    
    const t0 = performance.now();
    
    try {
        const fd = new FormData(); fd.append("file", f);
        const res = await fetch("/convert", { method: "POST", body: fd });
        const json = await res.json();
        
        const time = ((performance.now()-t0)/1000).toFixed(2) + "s";
        $("timerBadge").innerText = time;
        $("timerBadge").classList.remove("hidden");

        if(res.ok) {
            mdOut.value = json.markdown;
            updatePreview();
            
            // Prepare download
            if(currentBlobUrl) URL.revokeObjectURL(currentBlobUrl);
            const blob = new Blob([json.markdown], {type: "text/markdown"});
            currentBlobUrl = URL.createObjectURL(blob);
            
            $("downloadBtn").disabled = false;
            $("downloadBtn").onclick = () => {
                const a = document.createElement("a");
                a.href = currentBlobUrl;
                a.download = json.output_filename || "document.md";
                a.click();
            };
            
            $("statusBadge").innerText = "Terminé";
            $("statusBadge").style.color = "#10b981";
            showToast("Conversion terminée !");
        } else {
            mdOut.value = "Erreur: " + JSON.stringify(json, null, 2);
            showToast("Erreur serveur", "error");
        }
    } catch(e) {
        mdOut.value = e.toString();
        showToast("Erreur réseau", "error");
    } finally {
        btn.innerHTML = '<i class="fas fa-bolt"></i> Convertir';
        btn.disabled = false;
    }
}

function copyToClipboard() {
    if(!mdOut.value) return;
    navigator.clipboard.writeText(mdOut.value);
    showToast("Code Markdown copié !");
}

function downloadFile() {
    if(!currentBlobUrl) return;
    // Triggered by onclick re-assignment above
}

mdOut.addEventListener("input", updatePreview);

</script>
</body>
</html>
'''

@app.get("/", response_class=HTMLResponse)
def index(): return HTMLResponse(HTML_PAGE)

@app.get("/images/{conversion_id}/{filename}")
def get_converted_image(conversion_id: str, filename: str):
    safe_conversion = _safe_slug(conversion_id, "conversion")
    safe_filename = _safe_slug(filename, "image")
    path = os.path.abspath(os.path.join(IMAGE_OUTPUT_DIR, safe_conversion, safe_filename))
    root = os.path.abspath(IMAGE_OUTPUT_DIR)
    if not path.startswith(root + os.sep) or not os.path.exists(path):
        raise HTTPException(status_code=404, detail="Image not found")
    return FileResponse(path)

@app.post("/convert")
async def convert(file: UploadFile = File(...), inline_images: Optional[str] = Form(None)):
    try:
        t0 = time.perf_counter()
        content = await file.read()
        if SAVE_UPLOADS:
            with open(os.path.join(UPLOAD_DIR, file.filename), "wb") as f: f.write(content)

        ftype = guess_file_type(file.filename, file.content_type)
        conversion_id = uuid.uuid4().hex
        inline = _parse_bool(inline_images, MARKDOWN_INLINE_IMAGES)
        meta = {"type": ftype, "conversion_id": conversion_id, "inline_images": inline}
        images: List[Dict[str, Any]] = []
        
        if ftype == "pdf":
            md, images = render_pdf_markdown_inline(content, meta, file.filename, conversion_id, inline)
        elif ftype == "docx":
            md, m_docx, images = render_docx_markdown(content, file.filename, conversion_id, inline)
            meta.update(m_docx)
        elif ftype == "pptx":
            md, images = render_pptx_markdown(content, file.filename, conversion_id, inline)
        elif ftype == "json":
            md = render_json_markdown(content)
        elif ftype == "html":
            txt = content.decode("utf-8", errors="ignore")
            uris = _extract_html_data_imgs(txt)
            replacement_uris: List[str] = []
            for idx, uri in enumerate(uris, start=1):
                decoded = _data_uri_to_image(uri)
                if not decoded:
                    continue
                im, _mime = decoded
                asset = _save_image_asset(
                    im,
                    conversion_id=conversion_id,
                    source_filename=file.filename,
                    image_index=idx,
                    caption=f"{IMG_ALT_PREFIX} - image {idx}",
                )
                images.append(asset)
                replacement_uris.append(uri if inline else asset.get("public_url") or asset.get("relative_url") or "")
            eng = MarkItDown()
            res = eng.convert_stream(io.BytesIO(content), file_extension=".html")
            md = getattr(res, "text_content", "") or ""
            md = _html_img_datauri_to_markdown(md)
            md = _inject_full_data_uris_into_markdown(md, replacement_uris, IMG_ALT_PREFIX)
            if not inline:
                md = re.sub(r'!\[[^\]]*\]\([^)]+\)', "", md)
            md = _md_cleanup(md)
        elif ftype == "image":
            with Image.open(io.BytesIO(content)) as im:
                asset = _save_image_asset(
                    im,
                    conversion_id=conversion_id,
                    source_filename=file.filename,
                    image_index=1,
                    caption=file.filename,
                )
                images = [asset]
                if inline:
                    uri = _pil_to_base64(im, IMG_FORMAT, IMG_JPEG_QUALITY)
                    md = _image_markdown(asset, True, uri)
                else:
                    md = f"Image extraite: {file.filename}"
        else:
            eng = MarkItDown()
            try:
                res = eng.convert_stream(io.BytesIO(content), file_extension=os.path.splitext(file.filename)[1])
                md = getattr(res, "text_content", "") or ""
            except:
                md = content.decode("utf-8", errors="replace")

        out_name = f"{os.path.splitext(file.filename)[0]}.md"
        if SAVE_OUTPUTS:
            with open(os.path.join(OUTPUT_DIR, out_name), "w", encoding="utf-8") as f: f.write(md)

        meta["duration"] = round(time.perf_counter() - t0, 3)
        meta["images_count"] = len(images)
        return JSONResponse({
            "filename": file.filename,
            "output_filename": out_name,
            "markdown": md,
            "images": images,
            "metadata": meta,
        })

    except Exception as e:
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=str(e))
